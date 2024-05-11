from pptx import Presentation
import docx
import speech_recognition as sr
from PyPDF2 import PdfReader
import moviepy.editor as mp
import time
import requests
import uuid
import os


def current_millis():
    return round(time.time() * 1000)


class Extractor:
    def __init__(self, base_temp_dir, speech_recognition_model="base"):
        self.base_temp_dir = base_temp_dir
        self.speech_recognition_model = speech_recognition_model

    def extract_data(self, df, extractors):
        # id for downloading local file and assigning name to it if needed
        extraction_uuid = uuid.uuid4().hex
        initial_file_path = df.link
        # downloading file if needed
        if df.link.startswith("http"):
            initial_file_path = self.base_temp_dir + extraction_uuid + get_file_extension(df)
            download_file_from_url(df, initial_file_path)

        parameters = {
            "base_temp_dir": self.base_temp_dir,
            "initial_file_path": initial_file_path,
            "extraction_uuid": extraction_uuid,
            "extractors": extractors
        }
        result = df.extract_data(parameters)
        return result

    def remove_temp_files(self):
        to_cleanup = os.listdir(self.base_temp_dir)
        for file in to_cleanup:
            print(f"Removing {self.base_temp_dir}{file}")
            os.remove(self.base_temp_dir + ("" if self.base_temp_dir.endswith("/") else "/") + file)

    def extract_text_from_audio(self, path):
        r = sr.Recognizer()
        audio_file = sr.AudioFile(path)
        with audio_file as source:
            audio = r.record(source)
            return r.recognize_whisper(audio_data=audio, model=self.speech_recognition_model)


def download_file_from_url(df, dest_path):
    if df.auth_data is not None:
        auth_header = df.auth_data['type'].value['extract_auth_header'](df)
        response = requests.get(df.link, headers={"Authorization": auth_header})
    else:
        response = requests.get(df.link)
    print(f"Downloading content of digital footprint by link {df.link}")
    with open(dest_path, "wb") as file:
        file.write(response.content)


def extract_text_from_presentation(path):
    pres = Presentation(path)
    text = ""
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = text + "\n\n" + shape.text
    return text


def get_txt_from_pdf(pdf_path):
    fullText = ""
    # Load the PDF
    pdf_reader = PdfReader(pdf_path)
    for page in pdf_reader.pages:
        fullText = fullText + page.extract_text() + "\n\n"
    return fullText


def get_txt_from_doc(doc_path):
    doc = docx.Document(doc_path)
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    return '\n\n'.join(fullText)


def extract_text_from_document(path):
    if path.endswith(".docx") or path.endswith(".doc"):
        return get_txt_from_doc(path)
    elif path.endswith(".pdf"):
        return get_txt_from_pdf(path)
    elif path.endswith(".pptx"):
        return extract_text_from_presentation(path)
    raise f"Unrecognized file type {path}"


def get_file_extension(df, default_ext=""):
    if df.additional_info is not None and "file_name" in df.additional_info:
        return os.path.splitext(df.additional_info["file_name"])[1]
    return default_ext


def extract_audio_from_video(path, audio_path):
    # Load the video
    video = mp.VideoFileClip(path)
    # Extract the audio from the video
    audio_file = video.audio
    audio_path_result = audio_path + ".wav"
    audio_file.write_audiofile(audio_path_result)
    return audio_path_result
